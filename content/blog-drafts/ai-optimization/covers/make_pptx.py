#!/usr/bin/env python3
"""Generate editable .pptx covers (open in Canva) for the AI-visibility post."""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PX = 9525  # EMU per CSS px @96dpi
def px(v): return Emu(int(v * PX))
def ptpx(v): return Pt(v * 0.75)  # px -> points

CANVAS = RGBColor(0xfa, 0xf9, 0xf6)
INK    = RGBColor(0x35, 0x33, 0x2f)
YEL    = RGBColor(0xda, 0xff, 0x00)
GREY   = RGBColor(0x7a, 0x77, 0x6f)
DGREY  = RGBColor(0x47, 0x45, 0x3f)
HERE = os.path.dirname(os.path.abspath(__file__))

def deck(w, h):
    prs = Presentation()
    prs.slide_width = px(w); prs.slide_height = px(h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, px(w), px(h))
    bg.fill.solid(); bg.fill.fore_color.rgb = CANVAS; bg.line.fill.background()
    bg.shadow.inherit = False
    return prs, slide

def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1.2)
    s.shadow.inherit = False
    return s

def oval(slide, x, y, d, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(d), px(d))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1.4)
    s.shadow.inherit = False
    return s

def text(slide, x, y, w, h, runs, size, font="Futura", bold=False, upper=False,
         color=INK, align=PP_ALIGN.LEFT, spacing=None, leading=None):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    if leading: p.line_spacing = leading
    if isinstance(runs, str):
        runs = [(runs, bold, color)]
    for txt, b, c in runs:
        r = p.add_run(); r.text = txt.upper() if upper else txt
        f = r.font; f.size = ptpx(size); f.name = font; f.bold = b; f.color.rgb = c
        if spacing is not None:
            r._r.get_or_add_rPr().set("spc", str(int(spacing)))
    return tb

# ---------- BLOG 1200 x 630 ----------
prs, s = deck(1200, 630)
text(s, 78, 70, 1050, 30, "THE DOT CREATIVE     ·     MARKETING NOTES", 19,
     bold=True, upper=True, color=GREY, spacing=280)
rect(s, 300, 452, 96, 26, YEL)                       # highlight swipe behind "AI" (nudge in Canva)
text(s, 74, 168, 1060, 250, "Can AI find your business?", 94, color=INK, leading=1.0)
text(s, 78, 452, 900, 96,
     "The new kind of search you can't ignore, and how to check where you stand.", 28, color=DGREY)
for i in range(26):
    rect(s, 78 + i * 20, 556, 2, 13, INK)            # tick stripe
text(s, 78, 588, 700, 28, "thedotcreative.co", 20, bold=True, upper=True, color=INK, spacing=160)
blog_dots = [(0,0,'k'),(1,0,'o'),(2,0,'y'),(3,0,'o'),(0,1,'o'),(1,1,'y'),(2,1,'o'),(3,1,'k')]
bx, by, cell, gap = 968, 470, 30, 16
for cx, cy, kind in blog_dots:
    x = bx + cx * (cell + gap); y = by + cy * (cell + gap)
    if kind == 'k': oval(s, x, y, cell, fill=INK)
    elif kind == 'y': oval(s, x, y, cell, fill=YEL)
    else: oval(s, x, y, cell, fill=CANVAS, line=INK)
prs.save(os.path.join(HERE, "blog-cover.pptx"))

# ---------- IG 1080 x 1350 ----------
prs, s = deck(1080, 1350)
text(s, 84, 96, 900, 34, "THE DOT CREATIVE", 22, bold=True, upper=True, color=GREY, spacing=300)
rect(s, 442, 392, 132, 32, YEL)                      # highlight behind "AI"
text(s, 84, 150, 940, 320, "Can AI find your business?", 108, color=INK, leading=0.98)
for i in range(24):
    rect(s, 84 + i * 18, 486, 2, 14, INK)            # tick
rect(s, 120, 1052, 560, 52, YEL)                     # highlight under 45%
text(s, 80, 560, 940, 320, "45%", 280, color=INK, leading=0.9)
text(s, 84, 1120, 920, 220,
     [("of people now ask AI to find a local business, up from just ", False, DGREY),
      ("6% a year ago", True, INK),
      (". If it can't name you, you are not on the list.", False, DGREY)],
     34, color=DGREY, leading=1.3)
ig_dots = [(0,0,'k'),(1,0,'o'),(2,0,'y'),(0,1,'y'),(1,1,'k'),(2,1,'o')]
bx, by, cell, gap = 884, 1006, 30, 16
for cx, cy, kind in ig_dots:
    x = bx + cx * (cell + gap); y = by + cy * (cell + gap)
    if kind == 'k': oval(s, x, y, cell, fill=INK)
    elif kind == 'y': oval(s, x, y, cell, fill=YEL)
    else: oval(s, x, y, cell, fill=CANVAS, line=INK)
text(s, 84, 1250, 940, 40, "New on the blog, plus a free 5-minute self-check.", 30, bold=True, color=INK)
text(s, 84, 1300, 940, 30, "thedotcreative.co   ·   @thedotcreativeagency", 22,
     bold=True, upper=True, color=GREY, spacing=180)
prs.save(os.path.join(HERE, "ig-cover.pptx"))
print("saved blog-cover.pptx and ig-cover.pptx")
